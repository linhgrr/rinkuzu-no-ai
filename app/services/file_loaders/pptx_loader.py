"""
PowerPoint (PPTX) file loader for RAG pipeline

Uses python-pptx library to extract text content, slide notes,
and structural information from PowerPoint presentations.
"""
import io
from typing import List, BinaryIO, Dict, Any
from loguru import logger

from app.core.interfaces.file_loader import FileLoader, DocumentChunk


class PPTXLoader(FileLoader):
    """
    PowerPoint PPTX file loader for extracting educational content
    
    Features:
    - Text extraction from slides
    - Speaker notes extraction
    - Slide titles and content separation
    - Table and shape text extraction
    - Slide-based chunking with metadata
    """
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 100):
        """
        Initialize PPTX loader
        
        Args:
            chunk_size: Maximum characters per chunk
            chunk_overlap: Overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_extensions = ["pptx"]
    
    async def load(self, file_stream: BinaryIO, filename: str) -> List[DocumentChunk]:
        """
        Load and process PPTX file into chunks
        
        Args:
            file_stream: PPTX file binary stream
            filename: Original filename for metadata
            
        Returns:
            List of document chunks with slide content and metadata
        """
        try:
            from pptx import Presentation
            
            # Load presentation from stream
            if hasattr(file_stream, 'read'):
                file_stream.seek(0)  # Ensure we're at the beginning
            
            presentation = Presentation(file_stream)
            total_slides = len(presentation.slides)
            
            logger.info(f"📊 Processing PPTX: {filename} ({total_slides} slides)")
            
            all_chunks = []
            presentation_text = ""
            
            # Extract text from each slide
            for slide_idx, slide in enumerate(presentation.slides, 1):
                try:
                    slide_content = self._extract_slide_content(slide, slide_idx, total_slides, filename)
                    
                    if slide_content["full_text"].strip():
                        # Add slide content to full presentation text
                        presentation_text += f"\n--- Slide {slide_idx} ---\n{slide_content['full_text']}\n"
                        
                        # Create individual slide chunk if it's not too large
                        if len(slide_content["full_text"]) <= self.chunk_size:
                            chunk = DocumentChunk(
                                content=slide_content["full_text"],
                                metadata=slide_content["metadata"],
                                page_number=slide_idx,
                                chunk_index=len(all_chunks)
                            )
                            all_chunks.append(chunk)
                        
                        # Create separate chunks for titles, content, and notes if they're substantial
                        if slide_content["title"] and len(slide_content["title"]) > 20:
                            title_chunk = DocumentChunk(
                                content=slide_content["title"],
                                metadata={
                                    **slide_content["metadata"],
                                    "chunk_type": "slide_title",
                                    "content_type": "title"
                                },
                                page_number=slide_idx,
                                chunk_index=len(all_chunks)
                            )
                            all_chunks.append(title_chunk)
                        
                        if slide_content["notes"] and len(slide_content["notes"]) > 50:
                            notes_chunk = DocumentChunk(
                                content=slide_content["notes"],
                                metadata={
                                    **slide_content["metadata"],
                                    "chunk_type": "speaker_notes",
                                    "content_type": "notes"
                                },
                                page_number=slide_idx,
                                chunk_index=len(all_chunks)
                            )
                            all_chunks.append(notes_chunk)
                    
                except Exception as e:
                    logger.warning(f"Failed to extract content from slide {slide_idx} of {filename}: {e}")
                    continue
            
            # If presentation is large, create overlapping text chunks
            if len(presentation_text) > self.chunk_size:
                text_chunks = self._create_text_chunks(presentation_text)
                
                for i, chunk_text in enumerate(text_chunks):
                    chunk = DocumentChunk(
                        content=chunk_text.strip(),
                        metadata={
                            "filename": filename,
                            "total_slides": total_slides,
                            "source_type": "pptx",
                            "chunk_type": "presentation_text",
                            "chunk_size": len(chunk_text),
                            "content_type": "combined"
                        },
                        chunk_index=len(all_chunks) + i
                    )
                    all_chunks.append(chunk)
            
            logger.info(f"✅ PPTX processed: {filename} -> {len(all_chunks)} chunks")
            return all_chunks
            
        except ImportError:
            logger.error("python-pptx library not installed. Install with: pip install python-pptx")
            raise ValueError("PowerPoint processing library (python-pptx) not available")
        except Exception as e:
            logger.error(f"❌ Failed to process PPTX {filename}: {e}")
            raise ValueError(f"Failed to process PowerPoint file: {e}")
    
    def _extract_slide_content(self, slide, slide_idx: int, total_slides: int, filename: str) -> Dict[str, Any]:
        """
        Extract comprehensive content from a single slide
        
        Args:
            slide: python-pptx slide object
            slide_idx: Slide number (1-based)
            total_slides: Total number of slides
            filename: Presentation filename
            
        Returns:
            Dictionary with slide content and metadata
        """
        slide_title = ""
        slide_content = []
        slide_notes = ""
        
        # Extract title (usually the first text box or title placeholder)
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    # First substantial text is likely the title
                    if not slide_title and len(shape.text.strip()) < 200:
                        slide_title = shape.text.strip()
                    else:
                        slide_content.append(shape.text.strip())
                
                # Extract table content
                if hasattr(shape, 'table'):
                    table_text = self._extract_table_text(shape.table)
                    if table_text:
                        slide_content.append(f"[Bảng]\n{table_text}")
        
        # Extract speaker notes
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            try:
                if hasattr(slide.notes_slide, 'notes_text_frame'):
                    slide_notes = slide.notes_slide.notes_text_frame.text.strip()
            except:
                pass
        
        # Combine all slide content
        full_slide_text = []
        if slide_title:
            full_slide_text.append(f"Tiêu đề: {slide_title}")
        if slide_content:
            full_slide_text.append(f"Nội dung:\n" + "\n".join(slide_content))
        if slide_notes:
            full_slide_text.append(f"Ghi chú:\n{slide_notes}")
        
        full_text = "\n\n".join(full_slide_text)
        
        return {
            "title": slide_title,
            "content": "\n".join(slide_content),
            "notes": slide_notes,
            "full_text": full_text,
            "metadata": {
                "filename": filename,
                "slide_number": slide_idx,
                "total_slides": total_slides,
                "source_type": "pptx",
                "chunk_type": "slide",
                "has_title": bool(slide_title),
                "has_content": bool(slide_content),
                "has_notes": bool(slide_notes),
                "content_length": len(full_text)
            }
        }
    
    def _extract_table_text(self, table) -> str:
        """
        Extract text content from PowerPoint table
        
        Args:
            table: python-pptx table object
            
        Returns:
            Formatted table text
        """
        try:
            table_rows = []
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if hasattr(cell, 'text') else ""
                    row_cells.append(cell_text)
                if any(row_cells):  # Only add non-empty rows
                    table_rows.append(" | ".join(row_cells))
            
            return "\n".join(table_rows) if table_rows else ""
            
        except Exception as e:
            logger.warning(f"Failed to extract table content: {e}")
            return ""
    
    def supports_file_type(self, file_extension: str) -> bool:
        """Check if this loader supports the given file extension"""
        return file_extension.lower() in self.supported_extensions
    
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions"""
        return self.supported_extensions.copy()
    
    def _create_text_chunks(self, text: str) -> List[str]:
        """
        Create overlapping text chunks from large presentation text
        
        Args:
            text: Full presentation text to chunk
            
        Returns:
            List of text chunks with overlap
        """
        if len(text) <= self.chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + self.chunk_size
            
            # Try to break at slide boundary first
            if end < len(text):
                slide_break = text.find('\n--- Slide ', start, end)
                if slide_break > start + self.chunk_size // 2:
                    end = slide_break
                else:
                    # Fall back to word boundary
                    space_pos = text.rfind(' ', start, end)
                    if space_pos > start + self.chunk_size // 2:
                        end = space_pos
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = max(start + 1, end - self.chunk_overlap)
            
            # Prevent infinite loop
            if start >= len(text):
                break
        
        return chunks 